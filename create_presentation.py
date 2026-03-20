"""
APRICA PowerPoint Presentation Generator
Creates a professional 12-slide deck for team/boss presentation
Color palette: Forest Green + Gold + Cream (African-themed)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree
import copy

# ──────────────────────────────────────────────
# COLOR PALETTE
# ──────────────────────────────────────────────
FOREST    = RGBColor(0x1B, 0x43, 0x32)   # Deep forest green
FOREST_LT = RGBColor(0x2D, 0x6A, 0x4F)   # Lighter green
GOLD      = RGBColor(0xD4, 0xA0, 0x17)   # African gold
GOLD_LT   = RGBColor(0xF0, 0xC0, 0x4A)   # Light gold
TERRACOTTA= RGBColor(0xC6, 0x5D, 0x2E)   # Terracotta
CREAM     = RGBColor(0xFA, 0xF3, 0xE0)   # Warm cream
CREAM_DK  = RGBColor(0xEF, 0xE5, 0xCA)   # Darker cream
CHARCOAL  = RGBColor(0x1A, 0x1A, 0x1A)   # Near black
MID_GRAY  = RGBColor(0x6B, 0x5B, 0x45)   # Warm gray
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
SUCCESS   = RGBColor(0x27, 0xAE, 0x60)
DANGER    = RGBColor(0xE7, 0x4C, 0x3C)

# Slide dimensions: 16x9
W = Inches(13.33)
H = Inches(7.5)


def rgb_hex(r, g, b):
    return RGBColor(r, g, b)


def add_rect(slide, x, y, w, h, fill_color, line_color=None, line_width=0):
    from pptx.util import Pt
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, text, x, y, w, h, font_size, bold=False, color=CHARCOAL,
                 align=PP_ALIGN.LEFT, italic=False, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox


def add_feature_card(slide, x, y, w, h, icon, title, desc, accent_color=FOREST):
    """Add a feature card with icon, title, and description."""
    # Card background
    card = add_rect(slide, x, y, w, h, WHITE)
    # Left accent bar
    add_rect(slide, x, y, 0.07, h, accent_color)
    # Icon
    add_text_box(slide, icon, x+0.15, y+0.12, 0.5, 0.5, 22, align=PP_ALIGN.CENTER)
    # Title
    add_text_box(slide, title, x+0.15, y+0.55, w-0.25, 0.35, 13, bold=True, color=accent_color)
    # Description
    add_text_box(slide, desc, x+0.15, y+0.88, w-0.25, h-1.05, 10, color=MID_GRAY)


def add_stat_block(slide, x, y, number, label, color=FOREST):
    """Add a large stat display."""
    add_text_box(slide, number, x, y, 2.0, 0.8, 44, bold=True, color=color, align=PP_ALIGN.CENTER,
                 font_name="Georgia")
    add_text_box(slide, label, x, y+0.75, 2.0, 0.4, 10, color=MID_GRAY, align=PP_ALIGN.CENTER)


def kente_strip(slide, y_pos, strip_h=0.08):
    """Add a Kente-inspired color strip."""
    colors = [GOLD, FOREST, TERRACOTTA, GOLD, FOREST_LT, GOLD_LT, FOREST, TERRACOTTA, GOLD]
    seg_w = 13.33 / len(colors)
    for i, c in enumerate(colors):
        add_rect(slide, i * seg_w, y_pos, seg_w + 0.01, strip_h, c)


# ══════════════════════════════════════════════
# CREATE PRESENTATION
# ══════════════════════════════════════════════
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]  # Completely blank layout

# ──────────────────────────────────────────────
# SLIDE 1: TITLE SLIDE
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)

# Full dark background
add_rect(slide, 0, 0, 13.33, 7.5, FOREST)

# Kente strip at bottom
kente_strip(slide, 6.9, 0.6)

# Decorative circle accent
circle_bg = slide.shapes.add_shape(9, Inches(9.2), Inches(-0.5), Inches(5), Inches(5))
circle_bg.fill.solid()
circle_bg.fill.fore_color.rgb = FOREST_LT
circle_bg.line.fill.background()

# Title
add_text_box(slide, "APRICA", 0.8, 1.5, 8, 1.8, 80, bold=True, color=GOLD,
             font_name="Georgia", align=PP_ALIGN.LEFT)

# Tagline
add_text_box(slide, "African Language AI Platform", 0.8, 3.2, 8, 0.7, 26,
             color=CREAM, align=PP_ALIGN.LEFT, font_name="Calibri")

add_text_box(slide, "Transforming African Voices · Bridging Language Barriers",
             0.8, 3.85, 9, 0.5, 16, color=GOLD_LT, align=PP_ALIGN.LEFT)

# Supporting languages
add_text_box(slide, "Nigerian Pidgin  ·  Yoruba  ·  Hausa  ·  Swahili",
             0.8, 5.2, 10, 0.5, 15, color=CREAM, align=PP_ALIGN.LEFT, italic=True)

# Bottom info
add_text_box(slide, "Presented to Leadership Team  ·  March 2026",
             0.8, 6.2, 8, 0.4, 12, color=GOLD_LT, align=PP_ALIGN.LEFT)

# Globe emoji large (text-based)
add_text_box(slide, "🌍", 10.5, 1.8, 2, 2, 90, align=PP_ALIGN.CENTER)

# ──────────────────────────────────────────────
# SLIDE 2: EXECUTIVE SUMMARY
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, CREAM)
add_rect(slide, 0, 0, 13.33, 1.4, FOREST)
kente_strip(slide, 1.4, 0.07)

add_text_box(slide, "Executive Summary", 0.6, 0.3, 10, 0.8, 32, bold=True,
             color=GOLD, font_name="Georgia")

# Three pillars
pillars = [
    ("🌍", "THE PROBLEM", "Over 2,000 African languages exist yet are critically underrepresented in AI, digital content, and global business communications. Millions of Africans are excluded from the digital economy due to language barriers."),
    ("⚡", "THE SOLUTION", "APRICA is a comprehensive AI-powered platform that translates, transcribes, and synthesizes African languages with cultural authenticity — making African voices heard in the digital world."),
    ("🎯", "THE OPPORTUNITY", "A $47B global language services market, 1.4 billion Africans, growing smartphone penetration, and zero dominant AI platforms for African languages represent a historic first-mover opportunity."),
]

for i, (icon, title, text) in enumerate(pillars):
    x = 0.4 + i * 4.3
    add_rect(slide, x, 1.7, 4.1, 5.2, WHITE)
    add_rect(slide, x, 1.7, 4.1, 0.06, GOLD)
    add_text_box(slide, icon, x, 1.85, 4.1, 0.6, 32, align=PP_ALIGN.CENTER)
    add_text_box(slide, title, x, 2.45, 4.1, 0.45, 12, bold=True, color=FOREST,
                 align=PP_ALIGN.CENTER)
    add_rect(slide, x+1.5, 2.85, 1.1, 0.04, GOLD)
    add_text_box(slide, text, x+0.15, 2.95, 3.8, 3.8, 11, color=MID_GRAY)

# ──────────────────────────────────────────────
# SLIDE 3: PLATFORM OVERVIEW
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, CREAM)
add_rect(slide, 0, 0, 0.25, 7.5, FOREST)
add_rect(slide, 0, 0, 13.33, 0.08, GOLD)

add_text_box(slide, "Platform Architecture", 0.55, 0.25, 8, 0.7, 30, bold=True,
             color=FOREST, font_name="Georgia")
add_text_box(slide, "7 powerful AI-driven features in one unified platform",
             0.55, 0.9, 10, 0.4, 14, color=MID_GRAY)

# 7 feature boxes in a 4+3 grid
features = [
    ("📝", "Text Translation", "40M+ word pairs"),
    ("🌐", "Website Translator", "Full URL support"),
    ("📄", "Document Translator", "PDF + TXT + Summary"),
    ("🤖", "AI Chatbot", "Cultural AI companion"),
    ("🔊", "Text to Speech", "6 voice profiles"),
    ("🎙️", "Speech to Text", "Whisper AI engine"),
    ("🎬", "Video Translator", "Full dubbing pipeline"),
]

positions = [
    (0.4, 1.5), (3.65, 1.5), (6.9, 1.5), (10.15, 1.5),
    (2.02, 4.3), (5.27, 4.3), (8.52, 4.3),
]

for (icon, title, sub), (x, y) in zip(features, positions):
    add_rect(slide, x, y, 2.9, 2.5, WHITE)
    add_rect(slide, x, y, 2.9, 0.06, TERRACOTTA)
    add_text_box(slide, icon, x, y+0.15, 2.9, 0.6, 28, align=PP_ALIGN.CENTER)
    add_text_box(slide, title, x+0.1, y+0.75, 2.7, 0.45, 12, bold=True,
                 color=FOREST, align=PP_ALIGN.CENTER)
    add_text_box(slide, sub, x+0.1, y+1.2, 2.7, 0.5, 10,
                 color=MID_GRAY, align=PP_ALIGN.CENTER)

# ──────────────────────────────────────────────
# SLIDE 4: LANGUAGES & VOICES
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, FOREST)
kente_strip(slide, 0, 0.06)

add_text_box(slide, "Languages & Voice Profiles", 0.5, 0.25, 10, 0.7, 30, bold=True,
             color=GOLD, font_name="Georgia")
add_text_box(slide, "Authentic cultural voices across African, British & American accents",
             0.5, 0.88, 12, 0.4, 14, color=CREAM)

# Language cards
langs = [
    ("🇳🇬", "Nigerian Pidgin", "75M+ speakers", "West Africa's lingua franca"),
    ("🇳🇬", "Yoruba", "40M+ speakers", "SW Nigeria · Benin · Togo"),
    ("🇳🇬", "Hausa", "70M+ speakers", "N. Nigeria · Niger · Ghana"),
    ("🇰🇪", "Swahili", "200M+ speakers", "East & Central Africa"),
]

for i, (flag, lang, speakers, region) in enumerate(langs):
    x = 0.4 + i * 3.15
    add_rect(slide, x, 1.55, 2.85, 2.8, FOREST_LT)
    add_rect(slide, x, 1.55, 2.85, 0.07, GOLD)
    add_text_box(slide, flag, x, 1.7, 2.85, 0.6, 30, align=PP_ALIGN.CENTER)
    add_text_box(slide, lang, x+0.1, 2.35, 2.65, 0.45, 14, bold=True, color=GOLD,
                 align=PP_ALIGN.CENTER)
    add_text_box(slide, speakers, x+0.1, 2.8, 2.65, 0.35, 11, color=GOLD_LT,
                 align=PP_ALIGN.CENTER, bold=True)
    add_text_box(slide, region, x+0.1, 3.15, 2.65, 0.55, 9, color=CREAM,
                 align=PP_ALIGN.CENTER, italic=True)

# Voice profiles table
add_text_box(slide, "Voice Profiles Available", 0.5, 4.65, 8, 0.4, 14, bold=True, color=GOLD)

voices = [
    ("Amara", "African Female 🎤", "African accent, warm & natural"),
    ("Ade", "African Male 🎙️", "Deep Nigerian/West African tones"),
    ("Emma / James", "British Female & Male", "Professional BBC-style accents"),
    ("Zara / Marcus", "American Female & Male", "Contemporary US English accents"),
]

for i, (name, type_, desc) in enumerate(voices):
    x = 0.4 + (i % 2) * 6.4
    y = 5.15 + (i // 2) * 0.7
    add_rect(slide, x, y, 5.8, 0.6, RGBColor(0x1A, 0x3D, 0x2E))
    add_text_box(slide, f"{name}", x+0.2, y+0.08, 1.4, 0.44, 12, bold=True, color=GOLD)
    add_text_box(slide, f"{type_}  ·  {desc}", x+1.7, y+0.12, 3.9, 0.36, 10, color=CREAM)

# ──────────────────────────────────────────────
# SLIDE 5: KEY ADVANTAGES (PROS)
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, CREAM)
add_rect(slide, 0, 0, 13.33, 1.35, FOREST)
kente_strip(slide, 1.35, 0.07)

add_text_box(slide, "Key Advantages", 0.6, 0.22, 9, 0.65, 30, bold=True, color=GOLD, font_name="Georgia")
add_text_box(slide, "Why APRICA is a game-changer for African digital communications",
             0.6, 0.82, 11, 0.4, 13, color=CREAM)

pros = [
    ("🏆", "First-Mover Advantage", "No dominant AI platform exists for West African languages. APRICA fills a critical gap in the $47B language services market with authentic, culturally-aware translations."),
    ("⚡", "Comprehensive Feature Set", "7 integrated tools (text, URL, document, video, TTS, STT, chatbot) in one platform — eliminating the need for multiple fragmented tools and subscriptions."),
    ("🎤", "Authentic African Voices", "Multiple accent profiles including African, British, and American voices provide culturally resonant audio output — not robotic generic text-to-speech."),
    ("🤖", "GPT-4o Powered", "Leverages cutting-edge GPT-4o for culturally nuanced translations that understand idioms, proverbs, and cultural context — far beyond basic word-for-word translation."),
    ("🌍", "Massive Market Reach", "Targeting 385M+ speakers across 4 major African language groups, with potential to expand to all 2,000+ African languages as the platform scales."),
    ("📱", "Scalable Architecture", "Built on Flask/Python with clean API design. Easily deployable to cloud platforms (AWS, GCP, Azure) and can integrate with WhatsApp, Slack, and other platforms."),
]

for i, (icon, title, text) in enumerate(pros):
    row = i // 2
    col = i % 2
    x = 0.4 + col * 6.45
    y = 1.7 + row * 1.8

    add_rect(slide, x, y, 6.2, 1.65, WHITE)
    add_rect(slide, x, y, 0.07, 1.65, SUCCESS)

    add_text_box(slide, icon, x+0.15, y+0.15, 0.5, 0.45, 18)
    add_text_box(slide, title, x+0.75, y+0.15, 5.3, 0.4, 13, bold=True, color=FOREST)
    add_text_box(slide, text, x+0.15, y+0.65, 5.85, 0.85, 10, color=MID_GRAY)

# ──────────────────────────────────────────────
# SLIDE 6: MARKET OPPORTUNITY
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, CHARCOAL)
add_rect(slide, 0, 0, 4.5, 7.5, FOREST)
kente_strip(slide, 0, 0.07)
add_rect(slide, 4.5, 0, 0.04, 7.5, GOLD)

# Left column: Big stats
add_text_box(slide, "Market\nOpportunity", 0.4, 0.3, 3.8, 1.4, 26, bold=True, color=GOLD,
             font_name="Georgia")

stats = [
    ("$47B", "Global Language Services Market"),
    ("1.4B", "People in Africa"),
    ("2,000+", "African Languages Exist"),
    ("4", "Languages We Support Today"),
    ("385M+", "Target Speaker Population"),
]

for i, (num, label) in enumerate(stats):
    y = 1.8 + i * 1.0
    add_text_box(slide, num, 0.3, y, 2.3, 0.55, 30, bold=True, color=GOLD_LT,
                 font_name="Georgia", align=PP_ALIGN.RIGHT)
    add_text_box(slide, label, 2.7, y+0.08, 1.5, 0.45, 9, color=CREAM)

# Right column: Narrative
add_text_box(slide, "The African Language\nTechnology Gap", 4.8, 0.3, 8, 1.0, 22, bold=True,
             color=WHITE, font_name="Georgia")

points = [
    "95% of AI language models are trained primarily on English data",
    "Less than 0.5% of internet content is in African languages",
    "African businesses lose billions annually to translation costs",
    "Growing smartphone adoption creates urgent demand for local language AI",
    "International companies increasingly need African language content for expansion",
    "APRICA is positioned to become the foundational layer for African language AI",
]

for i, point in enumerate(points):
    y = 1.55 + i * 0.82
    add_rect(slide, 4.8, y, 0.05, 0.6, GOLD)
    add_text_box(slide, point, 5.1, y+0.08, 8.0, 0.55, 11.5, color=CREAM)

# Bottom stat banner
add_rect(slide, 4.8, 7.0, 8.3, 0.45, GOLD)
add_text_box(slide, "African internet users expected to reach 1 billion by 2030 — each needing local language content",
             4.9, 7.05, 8.1, 0.38, 10, bold=True, color=FOREST, align=PP_ALIGN.CENTER)

# ──────────────────────────────────────────────
# SLIDE 7: HOW IT WORKS (WORKFLOW)
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, CREAM)
add_rect(slide, 0, 0, 13.33, 1.3, FOREST)
kente_strip(slide, 1.3, 0.07)

add_text_box(slide, "How APRICA Works", 0.6, 0.22, 9, 0.65, 30, bold=True,
             color=GOLD, font_name="Georgia")
add_text_box(slide, "Seamless AI pipeline from input to authentic African language output",
             0.6, 0.82, 11, 0.4, 13, color=CREAM)

# Translation workflow
steps = [
    ("1", "INPUT", "User provides text, URL, document, audio, or video in any language"),
    ("2", "AI PROCESS", "GPT-4o analyzes context, cultural nuance, and linguistic patterns"),
    ("3", "TRANSLATE", "Content translated to target African language with cultural accuracy"),
    ("4", "VOICE SYNTH", "ElevenLabs generates authentic African voice output (optional)"),
    ("5", "OUTPUT", "Translated text, audio, or dubbed video delivered instantly"),
]

box_w = 2.2
for i, (num, title, desc) in enumerate(steps):
    x = 0.3 + i * 2.6
    y = 1.65

    # Circle number
    circ = slide.shapes.add_shape(9, Inches(x+0.85), Inches(y), Inches(0.5), Inches(0.5))
    circ.fill.solid()
    circ.fill.fore_color.rgb = FOREST
    circ.line.fill.background()
    add_text_box(slide, num, x+0.85, y, 0.5, 0.5, 14, bold=True, color=GOLD,
                 align=PP_ALIGN.CENTER)

    add_rect(slide, x, y+0.55, box_w, 2.8, WHITE)
    add_rect(slide, x, y+0.55, box_w, 0.06, GOLD)
    add_text_box(slide, title, x+0.05, y+0.7, box_w-0.1, 0.45, 11, bold=True,
                 color=FOREST, align=PP_ALIGN.CENTER)
    add_text_box(slide, desc, x+0.1, y+1.2, box_w-0.2, 2.0, 10, color=MID_GRAY,
                 align=PP_ALIGN.CENTER)

    # Arrow between steps
    if i < len(steps) - 1:
        add_text_box(slide, "→", x+box_w+0.05, y+1.6, 0.4, 0.4, 16, color=GOLD,
                     align=PP_ALIGN.CENTER, bold=True)

# Two sub-workflows
add_text_box(slide, "Video Translation Pipeline", 0.4, 5.0, 6, 0.4, 13, bold=True, color=FOREST)
video_steps = ["Extract Audio (FFmpeg)", "→", "Transcribe (Whisper)", "→",
               "Translate (GPT-4)", "→", "Voice (ElevenLabs)", "→", "Merge Video"]
x_pos = 0.4
for step in video_steps:
    if step == "→":
        add_text_box(slide, step, x_pos, 5.5, 0.3, 0.35, 11, color=GOLD, bold=True)
        x_pos += 0.3
    else:
        add_rect(slide, x_pos, 5.4, 1.4, 0.5, FOREST_LT)
        add_text_box(slide, step, x_pos+0.05, 5.48, 1.3, 0.35, 8.5, color=CREAM,
                     align=PP_ALIGN.CENTER)
        x_pos += 1.5

# ──────────────────────────────────────────────
# SLIDE 8: TECHNOLOGY STACK
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, CREAM)
add_rect(slide, 0, 0, 13.33, 1.3, CHARCOAL)
kente_strip(slide, 1.3, 0.07)

add_text_box(slide, "Technology Stack", 0.6, 0.22, 9, 0.65, 30, bold=True,
             color=GOLD, font_name="Georgia")
add_text_box(slide, "Enterprise-grade AI and cloud infrastructure",
             0.6, 0.82, 11, 0.4, 13, color=CREAM)

tech_layers = [
    {
        "label": "🧠  AI & MACHINE LEARNING",
        "color": FOREST,
        "items": [
            ("OpenAI GPT-4o", "Translation, Chatbot, Summarization"),
            ("OpenAI Whisper", "Speech-to-Text Transcription"),
            ("ElevenLabs", "Multilingual Voice Synthesis"),
            ("OpenAI TTS", "Fallback Text-to-Speech"),
        ]
    },
    {
        "label": "⚙️  BACKEND & APIs",
        "color": TERRACOTTA,
        "items": [
            ("Flask (Python)", "REST API backend framework"),
            ("Flask-CORS", "Cross-origin resource sharing"),
            ("PyPDF2", "PDF text extraction"),
            ("BeautifulSoup4", "Web scraping & URL translation"),
        ]
    },
    {
        "label": "🎬  MEDIA PROCESSING",
        "color": GOLD,
        "items": [
            ("FFmpeg", "Video/Audio extraction & merging"),
            ("MoviePy", "Video processing library"),
            ("PyDub", "Audio manipulation"),
            ("WebM/MP3/WAV", "Multi-format audio support"),
        ]
    },
]

for col, layer in enumerate(tech_layers):
    x = 0.4 + col * 4.3
    add_rect(slide, x, 1.6, 4.0, 0.45, layer["color"])
    add_text_box(slide, layer["label"], x+0.1, 1.65, 3.8, 0.38, 10, bold=True, color=WHITE)

    for i, (tech, desc) in enumerate(layer["items"]):
        y = 2.2 + i * 1.1
        add_rect(slide, x, y, 4.0, 1.0, WHITE)
        add_rect(slide, x, y, 0.06, 1.0, layer["color"])
        add_text_box(slide, tech, x+0.15, y+0.1, 3.7, 0.38, 12, bold=True, color=CHARCOAL)
        add_text_box(slide, desc, x+0.15, y+0.52, 3.7, 0.38, 10, color=MID_GRAY)

# ──────────────────────────────────────────────
# SLIDE 9: PROS & CONS ANALYSIS
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, CREAM)
add_rect(slide, 0, 0, 13.33, 1.3, FOREST)
kente_strip(slide, 1.3, 0.07)

add_text_box(slide, "Honest Assessment: Pros & Considerations",
             0.5, 0.22, 12, 0.65, 28, bold=True, color=GOLD, font_name="Georgia")
add_text_box(slide, "A transparent evaluation of strengths and areas for development",
             0.5, 0.82, 11, 0.4, 13, color=CREAM)

# PROS column (bigger, more prominent)
add_rect(slide, 0.3, 1.6, 8.5, 0.45, SUCCESS)
add_text_box(slide, "✅  ADVANTAGES  (Strong Position)", 0.4, 1.65, 8.2, 0.38, 12,
             bold=True, color=WHITE)

pros_list = [
    "First dedicated AI platform for West African languages — massive competitive moat",
    "7 integrated tools eliminate the need for multiple expensive single-purpose services",
    "GPT-4o ensures culturally nuanced translations far superior to basic word-for-word tools",
    "ElevenLabs integration provides authentic African accent voices (not generic robotic TTS)",
    "Modular API design allows easy integration with WhatsApp, Slack, websites, and mobile apps",
    "Addresses a documented $47B global market with virtually no African-language AI competition",
    "Video dubbing capability is unique — no comparable product exists for African languages",
]

for i, pro in enumerate(pros_list):
    y = 2.2 + i * 0.56
    add_rect(slide, 0.3, y, 8.5, 0.5, WHITE)
    add_rect(slide, 0.3, y, 0.05, 0.5, SUCCESS)
    add_text_box(slide, f"✓  {pro}", 0.5, y+0.08, 8.1, 0.38, 10, color=CHARCOAL)

# CONS column (smaller, honest)
add_rect(slide, 9.1, 1.6, 4.0, 0.45, TERRACOTTA)
add_text_box(slide, "⚠️  CONSIDERATIONS", 9.2, 1.65, 3.8, 0.38, 12, bold=True, color=WHITE)

cons_list = [
    ("API Costs", "OpenAI/ElevenLabs usage fees scale with volume — needs pricing model"),
    ("Language Depth", "4 languages now; expanding to all African languages takes time"),
    ("Data Privacy", "Audio/document uploads require clear data policies"),
    ("Offline Mode", "Currently requires internet; offline capability is future roadmap"),
]

for i, (title, desc) in enumerate(cons_list):
    y = 2.2 + i * 1.2
    add_rect(slide, 9.1, y, 4.0, 1.1, WHITE)
    add_rect(slide, 9.1, y, 4.0, 0.06, TERRACOTTA)
    add_text_box(slide, title, 9.2, y+0.1, 3.8, 0.38, 11, bold=True, color=TERRACOTTA)
    add_text_box(slide, desc, 9.2, y+0.5, 3.8, 0.52, 10, color=MID_GRAY)

# Bottom mitigation note
add_rect(slide, 9.1, 7.05, 4.0, 0.4, CREAM_DK)
add_text_box(slide, "All considerations have clear mitigation paths →",
             9.2, 7.1, 3.8, 0.3, 9, color=MID_GRAY, italic=True)

# ──────────────────────────────────────────────
# SLIDE 10: USE CASES & IMPACT
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, CHARCOAL)
kente_strip(slide, 0, 0.07)
add_rect(slide, 0, 0.07, 13.33, 1.25, RGBColor(0x22, 0x22, 0x22))

add_text_box(slide, "Use Cases & Impact Sectors", 0.5, 0.18, 11, 0.7, 30, bold=True,
             color=GOLD, font_name="Georgia")
add_text_box(slide, "Where APRICA creates real-world value",
             0.5, 0.82, 10, 0.4, 13, color=CREAM)

use_cases = [
    ("🏢", "Business", "Corporate communications, HR documents, contracts, and training materials translated for African markets"),
    ("📰", "Media & News", "News articles, broadcasts, and social media content made accessible in local languages"),
    ("🏥", "Healthcare", "Medical information, patient forms, and health campaigns in native languages save lives"),
    ("🎓", "Education", "Educational content, textbooks, and e-learning platforms localized for African students"),
    ("🏛️", "Government", "Public policy documents, civic information, and official communications for all citizens"),
    ("🎵", "Entertainment", "Film dubbing, music lyrics, podcasts, and content creation in African languages"),
]

for i, (icon, sector, desc) in enumerate(use_cases):
    row = i // 3
    col = i % 3
    x = 0.4 + col * 4.3
    y = 1.55 + row * 2.75

    add_rect(slide, x, y, 4.0, 2.5, FOREST)
    add_rect(slide, x, y, 4.0, 0.07, GOLD)
    add_text_box(slide, icon, x, y+0.12, 4.0, 0.55, 28, align=PP_ALIGN.CENTER)
    add_text_box(slide, sector, x+0.1, y+0.72, 3.8, 0.38, 14, bold=True, color=GOLD,
                 align=PP_ALIGN.CENTER)
    add_text_box(slide, desc, x+0.15, y+1.1, 3.7, 1.25, 10, color=CREAM)

# ──────────────────────────────────────────────
# SLIDE 11: ROADMAP
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, CREAM)
add_rect(slide, 0, 0, 13.33, 1.3, FOREST)
kente_strip(slide, 1.3, 0.07)

add_text_box(slide, "Product Roadmap", 0.6, 0.22, 9, 0.65, 30, bold=True,
             color=GOLD, font_name="Georgia")
add_text_box(slide, "Planned milestones for growth and feature expansion",
             0.6, 0.82, 11, 0.4, 13, color=CREAM)

phases = [
    {
        "phase": "PHASE 1", "timeline": "Months 1-3",
        "title": "Foundation", "color": FOREST,
        "items": ["✓ 4 core African languages", "✓ Text & URL translation",
                  "✓ Document translation", "✓ Basic TTS/STT", "✓ AI Chatbot"]
    },
    {
        "phase": "PHASE 2", "timeline": "Months 4-6",
        "title": "Enhancement", "color": TERRACOTTA,
        "items": ["→ Mobile app (iOS/Android)", "→ 8 additional languages",
                  "→ Real-time voice chat", "→ API marketplace", "→ Enterprise dashboard"]
    },
    {
        "phase": "PHASE 3", "timeline": "Months 7-12",
        "title": "Scale", "color": GOLD,
        "items": ["⬡ WhatsApp/Telegram bot", "⬡ 20+ African languages",
                  "⬡ Offline capability", "⬡ Custom voice training", "⬡ B2B white-labeling"]
    },
    {
        "phase": "PHASE 4", "timeline": "Year 2+",
        "title": "Expansion", "color": MID_GRAY,
        "items": ["⬡ All 2,000+ African languages", "⬡ AI language learning platform",
                  "⬡ Live interpretation", "⬡ Sign language support", "⬡ Pan-African data center"]
    },
]

for i, phase in enumerate(phases):
    x = 0.35 + i * 3.18
    add_rect(slide, x, 1.6, 3.0, 0.7, phase["color"])
    add_text_box(slide, phase["phase"], x+0.1, 1.65, 1.5, 0.32, 10, bold=True, color=WHITE)
    add_text_box(slide, phase["timeline"], x+0.1, 1.95, 2.8, 0.28, 9, color=WHITE, italic=True)
    add_text_box(slide, phase["title"], x+0.1, 2.32, 2.8, 0.42, 14, bold=True, color=CHARCOAL)

    for j, item in enumerate(phase["items"]):
        y = 2.8 + j * 0.72
        add_rect(slide, x, y, 3.0, 0.65, WHITE)
        add_text_box(slide, item, x+0.15, y+0.12, 2.7, 0.42, 10.5, color=CHARCOAL)

# Current position indicator
add_rect(slide, 0, 6.8, 13.33, 0.08, GOLD)
add_text_box(slide, "◀ YOU ARE HERE  (Phase 1 Complete — Ready for Phase 2 Investment)",
             0.4, 6.88, 12.5, 0.45, 11, bold=True, color=FOREST, align=PP_ALIGN.CENTER)

# ──────────────────────────────────────────────
# SLIDE 12: CALL TO ACTION / CLOSING
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, FOREST)
kente_strip(slide, 0, 0.08)
kente_strip(slide, 7.42, 0.08)

# Large decorative element
globe_bg = slide.shapes.add_shape(9, Inches(7.5), Inches(0.5), Inches(7), Inches(7))
globe_bg.fill.solid()
globe_bg.fill.fore_color.rgb = FOREST_LT
globe_bg.line.fill.background()

add_text_box(slide, "🌍", 9.0, 1.0, 5, 3.5, 130, align=PP_ALIGN.CENTER)

# Left content
add_text_box(slide, "APRICA", 0.6, 0.5, 7, 1.2, 72, bold=True, color=GOLD,
             font_name="Georgia")
add_text_box(slide, "The Future of African\nLanguage Technology", 0.6, 1.65, 8, 1.2, 22,
             color=CREAM, font_name="Calibri")

add_text_box(slide, "Ready to launch. Ready to scale. Ready to impact.", 0.6, 3.0, 8, 0.5,
             16, color=GOLD_LT, italic=True)

# Next steps
add_text_box(slide, "NEXT STEPS", 0.6, 3.7, 3.5, 0.4, 12, bold=True, color=GOLD,
             font_name="Calibri")

next_steps = [
    "1. Live platform demo (10 minutes)",
    "2. API key configuration & testing",
    "3. Define pilot user group (recommended: 50 users)",
    "4. Budget approval for API costs (Phase 2)",
]

for i, step in enumerate(next_steps):
    add_text_box(slide, step, 0.6, 4.15 + i * 0.58, 7.8, 0.5, 13, color=CREAM)

# Bottom signature
add_rect(slide, 0, 6.8, 8.5, 0.6, RGBColor(0x12, 0x30, 0x24))
add_text_box(slide, "APRICA Platform  ·  v1.0.0  ·  Built for the African Digital Economy",
             0.3, 6.88, 8.0, 0.42, 11, color=GOLD_LT, align=PP_ALIGN.CENTER)

# ──────────────────────────────────────────────
# SAVE
# ──────────────────────────────────────────────
output_path = "/home/claude/aprica/APRICA_Platform_Presentation.pptx"
prs.save(output_path)
print(f"✅ Presentation saved: {output_path}")
print(f"   Slides: {len(prs.slides)}")
